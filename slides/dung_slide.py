# -*- coding: utf-8 -*-
"""Dựng bộ slide bảo vệ khóa luận: 24 slide chính + 8 slide phụ lục.

    python slides/noi_dung.py

Quy ước thiết kế (chốt sau vòng deep-research):
  - Dùng đúng khung trang trí của Template.pptx (HUFLIT), khổ 20 x 11,25 inch
  - Tiêu đề mỗi slide nội dung là MỘT CÂU KHẲNG ĐỊNH nêu kết luận
    (phương pháp assertion-evidence, Alley — Penn State)
  - Sàn cỡ chữ 30pt thật cho phần thân, theo hướng dẫn CVPR
  - Dải điều hướng chương nằm trong thanh tiêu đề của mẫu

Về hệ toạ độ: nội dung được soạn trên hệ ảo 13,333 x 7,5 inch, còn mẫu là
20 x 11,25 inch. Hai khổ chênh nhau đúng 1,5 lần nên mọi kích thước (toạ độ
và cỡ chữ) chỉ cần nhân K = 1,5 lúc vẽ. Nhờ vậy phần nội dung giữ nguyên,
không phải tính lại từng con số.
"""
import os
import sys

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

import khung_template as kt

GOC = os.path.dirname(os.path.abspath(__file__))
HINH = os.path.join(GOC, "hinh")
RA = os.path.join(GOC, "KLTN - Bao ve - Danh Hoang Hieu Nghi - K29.pptx")

# Hệ ảo 13,333 x 7,5 -> khổ mẫu 20 x 11,25
K = 1.5


def IN(v):
    """Inch trên hệ ảo -> Emu trên khổ thật."""
    return Inches(v * K)


def PT(v):
    """Cỡ chữ trên hệ ảo -> cỡ chữ thật."""
    return Pt(v * K)

# ── màu ────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x4E, 0x79)
XANH = RGBColor(0x2E, 0x75, 0xB6)
THAN = RGBColor(0x2D, 0x2D, 0x2D)
XAM = RGBColor(0x77, 0x77, 0x77)
KE = RGBColor(0xCC, 0xCC, 0xCC)
VANG = RGBColor(0xFF, 0xF2, 0xCC)
VIEN_VANG = RGBColor(0xE6, 0xC8, 0x00)
NAU = RGBColor(0x7A, 0x52, 0x00)
DO = RGBColor(0xC0, 0x00, 0x00)
TRANG = RGBColor(0xFF, 0xFF, 0xFF)

# Màu lấy từ Template.pptx của HUFLIT
HUFLIT = RGBColor(0x18, 0x45, 0x9E)   # xanh tiêu đề của mẫu
DO_MAU = RGBColor(0xE0, 0x23, 0x24)   # đỏ nhấn tên đề tài
NHAT = RGBColor(0xF0, 0xF4, 0xF8)
XANH_NHAT = RGBColor(0xEB, 0xF3, 0xFA)
SANG = RGBColor(0xA0, 0xBB, 0xDD)

FONT = "Arial"

# ── toạ độ (hệ ảo 13,333 x 7,5; nhân K = 1,5 khi vẽ) ───────────────────
W, H = 13.333, 7.5
LE = 0.6
RONG = W - 2 * LE

# Khung mẫu chiếm sẵn hai dải: thanh tiêu đề trên (hết ở 0,59) và dải chân
# trang dưới (bắt đầu từ 6,99). Vùng nội dung nằm gọn giữa hai mốc này.
Y_MAU_TREN = 0.60
Y_MAU_DUOI = 6.98

Y_BC, BC_H = 0.66, 0.26          # dải điều hướng chương
Y_TD, H_TD = 0.99, 0.84          # tiêu đề khẳng định
Y_GACH = 1.87                    # gạch ngang dưới tiêu đề
Y_ND = 2.03                      # mốc trên của vùng nội dung
Y_NGUON = 6.26                   # mốc trên của khối trích nguồn
H_NGUON = 0.68
CAO_ND = Y_NGUON - Y_ND          # chiều cao dùng được cho nội dung

CHUONG = ["Ch.1  Mở đầu", "Ch.2  Cơ sở lý thuyết", "Ch.3  Phương pháp",
          "Ch.4  Thực nghiệm", "Ch.5  Kết luận"]


# ── tiện ích ───────────────────────────────────────────────────────────
def hop_chu(s, x, y, w, h, canh=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = canh
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def dat(p, doan, cỡ=24, mau=THAN, dam=False, nghieng=False, gian=1.28,
        truoc=0, sau=8, canh=PP_ALIGN.LEFT):
    """doan: chuỗi, hoặc list các (text, dam, mau, cỡ) — cỡ/mau None = kế thừa."""
    p.alignment = canh
    p.line_spacing = gian
    p.space_before = Pt(truoc)
    p.space_after = Pt(sau)
    phan = [(doan, dam, mau, cỡ)] if isinstance(doan, str) else doan
    for txt, d, m, c in phan:
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = PT(c if c else cỡ)
        r.font.bold = bool(d)
        r.font.italic = nghieng
        r.font.color.rgb = m if m else mau
    return p


def khoi(tf, dong, cỡ=24, gian=1.28, sau=10):
    """dong: list các đoạn; mỗi đoạn là chuỗi hoặc list run."""
    for i, d in enumerate(dong):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        dat(p, d, cỡ=cỡ, gian=gian, sau=sau)


def gach_dau(txt, nhan=None, cỡ=24):
    """Một dòng gạch đầu dòng, có thể kèm nhãn in đậm ở đầu."""
    r = [("▪  ", False, XANH, cỡ)]
    if nhan:
        r.append((nhan, True, THAN, cỡ))
    r.append((txt, False, THAN, cỡ))
    return r


def chu_nhat(s, x, y, w, h, to=None, vien=None, day=1.0, bo_tron=False):
    hd = MSO_SHAPE.ROUNDED_RECTANGLE if bo_tron else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(hd, IN(x), IN(y), IN(w), IN(h))
    if bo_tron:
        sp.adjustments[0] = 0.06
    if to is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = to
    if vien is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = vien
        sp.line.width = PT(day)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp


def nen(s, mau=TRANG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = mau


def breadcrumb(s, ch):
    """Dải định vị chương, đặt ngay dưới thanh tiêu đề của mẫu."""
    if not ch:
        return
    bw = RONG / 5
    for i, ten in enumerate(CHUONG):
        hoat = (i + 1) == ch
        tf = hop_chu(s, LE + i * bw, Y_BC, bw, BC_H, canh=MSO_ANCHOR.MIDDLE)
        dat(tf.paragraphs[0], ten, cỡ=11, mau=NAVY if hoat else XAM,
            dam=hoat, sau=0, canh=PP_ALIGN.CENTER)
        if hoat:
            chu_nhat(s, LE + i * bw + bw * 0.10, Y_BC + BC_H - 0.04,
                     bw * 0.80, 0.035, to=DO)


# Tên mục hiển thị trên thanh tiêu đề của mẫu, theo từng chương
TEN_MUC = {
    1: "MỞ ĐẦU",
    2: "TỔNG QUAN CƠ SỞ LÝ THUYẾT",
    3: "PHƯƠNG PHÁP ĐỀ XUẤT",
    4: "THỰC NGHIỆM VÀ ĐÁNH GIÁ",
    5: "KẾT LUẬN",
}


def tieu_de(s, txt, ch=None, nho=False):
    """Tiêu đề khẳng định của slide, nằm dưới thanh tiêu đề mẫu."""
    if ch:
        nhan_muc(s, "%02d" % ch, TEN_MUC[ch])
    breadcrumb(s, ch)
    cỡ = 26
    if len(txt) > 96:
        cỡ = 22
    elif len(txt) > 66:
        cỡ = 24
    if nho:
        cỡ -= 2
    tf = hop_chu(s, LE, Y_TD, RONG, H_TD)
    dat(tf.paragraphs[0], txt, cỡ=cỡ, mau=NAVY, dam=True, gian=1.14, sau=0)
    chu_nhat(s, LE, Y_GACH, RONG, 0.018, to=KE)


# Danh mục tài liệu tham khảo — số hiệu GIỮ ĐÚNG như trong khóa luận,
# để hội đồng đối chiếu chéo giữa slide và bài viết không bị lệch.
TLTK = {
    1: "[1] G. N. Hounsfield, “Computerized transverse axial scanning "
       "(tomography): Part 1,” Br. J. Radiol., 1973",
    2: "[2] G. Litjens và cộng sự, “A survey on deep learning in medical "
       "image analysis,” Med. Image Anal., 2017",
    3: "[3] H. R. Roth và cộng sự, “DeepOrgan: Multi-level Deep "
       "Convolutional Networks for Automated Pancreas Segmentation,” "
       "MICCAI, 2015",
    4: "[4] J. Ma và cộng sự, “AbdomenCT-1K: Is Abdominal Organ Segmentation "
       "a Solved Problem?,” IEEE TPAMI, 2022",
    5: "[5] H. Sung và cộng sự, “Global Cancer Statistics 2020: GLOBOCAN,” "
       "CA Cancer J. Clin., 2021",
    6: "[6] O. Ronneberger, P. Fischer, T. Brox, “U-Net: Convolutional "
       "Networks for Biomedical Image Segmentation,” MICCAI, 2015",
    7: "[7] A. Vaswani và cộng sự, “Attention Is All You Need,” NeurIPS, 2017",
    8: "[8] J. Chen và cộng sự, “TransUNet: Transformers Make Strong Encoders "
       "for Medical Image Segmentation,” arXiv:2102.04306, 2021",
    9: "[9] S. Woo, J. Park, J.-Y. Lee, I. S. Kweon, “CBAM: Convolutional "
       "Block Attention Module,” ECCV, 2018",
    10: "[10] B. Landman và cộng sự, “MICCAI Multi-Atlas Labeling Beyond the "
        "Cranial Vault – Workshop and Challenge,” Synapse, 2015",
    11: "[11] K. He, X. Zhang, S. Ren, J. Sun, “Deep Residual Learning for "
        "Image Recognition,” CVPR, 2016",
    12: "[12] A. Dosovitskiy và cộng sự, “An Image Is Worth 16x16 Words,” "
        "ICLR, 2021",
    14: "[14] J. Long, E. Shelhamer, T. Darrell, “Fully Convolutional "
        "Networks for Semantic Segmentation,” CVPR, 2015",
    15: "[15] F. Milletari, N. Navab, S.-A. Ahmadi, “V-Net: Fully "
        "Convolutional Neural Networks for Volumetric Medical Image "
        "Segmentation,” 3DV, 2016",
    16: "[16] A. A. Taha, A. Hanbury, “Metrics for evaluating 3D medical "
        "image segmentation,” BMC Medical Imaging, 2015",
    17: "[17] Z. Zhou và cộng sự, “UNet++: A Nested U-Net Architecture for "
        "Medical Image Segmentation,” DLMIA, 2018",
    19: "[19] O. Oktay và cộng sự, “Attention U-Net: Learning Where to Look "
        "for the Pancreas,” arXiv:1804.03999, 2018",
    22: "[22] H. Cao và cộng sự, “Swin-Unet: Unet-Like Pure Transformer for "
        "Medical Image Segmentation,” ECCV Workshops, 2023",
    23: "[23] H. Wang, P. Cao, J. Wang, O. R. Zaiane, “UCTransNet: Rethinking "
        "the Skip Connections in U-Net,” AAAI, 2022",
    25: "[25] J. Hu, L. Shen, G. Sun, “Squeeze-and-Excitation Networks,” "
        "CVPR, 2018",
    26: "[26] F. Lu, J. Xu, Q. Sun, Q. Lou, “FMD-TransUNet,” Applied "
        "Intelligence, 2025",
    27: "[27] J. Chen và cộng sự, “TransUNet: Rethinking the U-Net "
        "Architecture Design,” Medical Image Analysis, 2024",
    28: "[28] S. Chen, X. Tan, B. Wang, X. Hu, “Reverse Attention for Salient "
        "Object Detection,” ECCV, 2018",
    29: "[29] D.-P. Fan và cộng sự, “PraNet: Parallel Reverse Attention "
        "Network for Polyp Segmentation,” MICCAI, 2020",
}


def nguon(s, refs=(), tu_do=None):
    """Chân slide: liệt kê nguồn đầy đủ cho mọi luận điểm có trích dẫn.

    refs   = danh sách số hiệu TLTK (đúng như trong khóa luận)
    tu_do  = phần ghi thêm cho số liệu do khóa luận tự đo
    """
    dong = [TLTK[n] for n in refs]
    if tu_do:
        dong.append(tu_do)
    if not dong:
        return
    # neo đáy: một nguồn nằm sát chân slide, nhiều nguồn nở ngược lên trên
    tf = hop_chu(s, LE, Y_NGUON, RONG, H_NGUON, canh=MSO_ANCHOR.BOTTOM)
    for i, d in enumerate(dong):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        dat(p, d, cỡ=10.5, mau=XAM, gian=1.04, sau=0)


def nhan_cot(s, txt, x, y, w):
    tf = hop_chu(s, x, y, w, 0.36)
    dat(tf.paragraphs[0], txt, cỡ=21, mau=XANH, dam=True, sau=0)


def hinh(s, ten, x, y, w_max, h_max):
    """Đặt ảnh vừa khít khung, giữ đúng tỷ lệ, canh giữa khung."""
    p = os.path.join(HINH, ten)
    iw, ih = Image.open(p).size
    ti = iw / ih
    w, h = w_max, w_max / ti
    if h > h_max:
        h, w = h_max, h_max * ti
    s.shapes.add_picture(p, IN(x + (w_max - w) / 2),
                         IN(y + (h_max - h) / 2), IN(w), IN(h))


def bang(s, dl, x, y, w, cot_w=None, cỡ=17, cao=0.42, dam_dong=None,
         canh_giua_tu=1, do_dong=()):
    """dl[0] là hàng tiêu đề.

    dam_dong = chỉ số hàng thân được tô nổi (vàng + chữ navy đậm)
    do_dong  = các chỉ số hàng cần tô chữ đỏ (giá trị bất lợi)
    """
    nr, nc = len(dl), len(dl[0])
    sh = s.shapes.add_table(nr, nc, IN(x), IN(y), IN(w), IN(cao * nr))
    t = sh.table
    t.first_row = True
    t.horz_banding = False
    if cot_w:
        tong = sum(cot_w)
        for i, cw in enumerate(cot_w):
            t.columns[i].width = Emu(int(IN(w) * cw / tong))
    for r in range(nr):
        t.rows[r].height = IN(cao)
        for c in range(nc):
            o = t.cell(r, c)
            o.text = ""
            o.margin_left = o.margin_right = IN(0.09)
            o.margin_top = o.margin_bottom = IN(0.03)
            o.vertical_anchor = MSO_ANCHOR.MIDDLE
            o.fill.solid()
            if r == 0:
                o.fill.fore_color.rgb = NAVY
            elif dam_dong is not None and r == dam_dong:
                o.fill.fore_color.rgb = VANG
            else:
                o.fill.fore_color.rgb = TRANG if r % 2 else NHAT
            p = o.text_frame.paragraphs[0]
            noi_bat = dam_dong is not None and r == dam_dong
            if r == 0:
                mau = TRANG
            elif noi_bat:
                mau = NAVY
            elif r in do_dong:
                mau = DO
            else:
                mau = THAN
            dat(p, str(dl[r][c]), cỡ=cỡ, mau=mau,
                dam=(r == 0 or noi_bat or r in do_dong), sau=0, gian=1.05,
                canh=PP_ALIGN.LEFT if c < canh_giua_tu else PP_ALIGN.CENTER)
    return t


def hop_nhan_manh(s, x, y, w, h, dong, cỡ=25):
    chu_nhat(s, x, y, w, h, to=XANH_NHAT, vien=XANH, day=1.5, bo_tron=True)
    tf = hop_chu(s, x + 0.3, y + 0.2, w - 0.6, h - 0.4,
                 canh=MSO_ANCHOR.MIDDLE)
    khoi(tf, dong, cỡ=cỡ, gian=1.26, sau=6)


def nhan_phu_luc(s, txt):
    tf = hop_chu(s, LE, Y_BC, RONG, BC_H)
    dat(tf.paragraphs[0], txt, cỡ=12, mau=XAM, nghieng=True, sau=0)


# ── khung của Template.pptx ────────────────────────────────────────────
# Mở mẫu một lần, dùng chính nó làm tệp đầu ra để giữ nguyên theme, master
# và layout. Các slide mẫu được xóa sau khi đã sao chép xong phần trang trí.
_pr_mau = kt.mo_mau()
_mau_bia = _pr_mau.slides[kt.BIA]
_mau_muc_luc = _pr_mau.slides[kt.MUC_LUC]
_mau_chuong = _pr_mau.slides[kt.CHUONG]
_mau_nd = _pr_mau.slides[kt.NOI_DUNG]
_mau_cam_on = _pr_mau.slides[kt.CAM_ON]

# Trang nội dung giữ nguyên toàn bộ khung, kể cả hai group chữ trên thanh
# tiêu đề: "Group 6" là số chương, "Group 9" là tên mục. Hai group này được
# giữ lại và chỉ thay nội dung chữ (xem nhan_muc), nhờ vậy font, cỡ, màu và
# vị trí vẫn đúng như mẫu.
_BO_ND = ()
_G_SO_CHUONG = "Group 6"
_G_TEN_MUC = "Group 9"
_BO_BIA = ("PHÁT TRIỂN KỸ THUẬT", "KHỐI U TRÊN DA", "Giảng viên hướng dẫn",
           "KHOÁ LUẬN TỐT NGHIỆP", "Khoa Công Nghệ Thông Tin",
           "CHUYÊN NGÀNH", "TP.HCM, tháng", "TRƯỜNG ĐẠI HỌC NGOẠI NGỮ",
           "HO CHI MINH CITY UNIVERSITY")
_BO_CAM_ON = ("TRÂN TRỌNG CẢM ƠN",)


def bo_phong():
    """Trả về Presentation dùng để dựng bài, đã mang theme của mẫu."""
    return _pr_mau


def moi(pr, ch=None, khung="nd"):
    """Slide mới trên nền Blank, có sẵn khung trang trí của mẫu.

    khung = 'nd' trang nội dung | 'bia' | 'chuong' | 'cam_on' | None
    """
    s = pr.slides.add_slide(pr.slide_masters[0].slide_layouts[6])
    nen(s)
    if khung == "nd":
        kt.sao_khung(_mau_nd, s, bo_chu=_BO_ND)
    elif khung == "bia":
        kt.sao_khung(_mau_bia, s, bo_chu=_BO_BIA)
    elif khung == "chuong":
        kt.sao_khung(_mau_nd, s, bo_chu=_BO_ND)
    elif khung == "cam_on":
        kt.sao_khung(_mau_cam_on, s, bo_chu=_BO_CAM_ON)
    return s


def nhan_muc(s, so, ten):
    """Điền số chương và tên mục vào thanh tiêu đề của mẫu."""
    kt.dat_chu_theo_ten(s, _G_SO_CHUONG, so)
    kt.dat_chu_theo_ten(s, _G_TEN_MUC, ten)


# ── toạ độ thật của khổ mẫu (20 x 11,25 inch) ──────────────────────────
# Trang bìa và trang cảm ơn được dựng theo đúng toạ độ của Template.pptx,
# nên hai hàm dưới nhận inch thật rồi tự quy về hệ ảo.
def hop_that(s, x, y, w, h, canh=MSO_ANCHOR.TOP):
    return hop_chu(s, x / K, y / K, w / K, h / K, canh)


def ct(v):
    """Cỡ chữ thật -> cỡ chữ trên hệ ảo."""
    return v / K
