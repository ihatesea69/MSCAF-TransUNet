# -*- coding: utf-8 -*-
"""Lấy khung trang trí (chrome) từ Template.pptx và dán sang slide mới.

Template.pptx là bộ slide mẫu của HUFLIT, khổ 20 x 11,25 inch. Bài này giữ
nguyên phần trang trí của mẫu (thanh tiêu đề, số chương, dải chân trang) rồi
đổ nội dung khóa luận vào giữa.

Khó khăn kỹ thuật: một số hình trang trí là freeform có ảnh nhúng (png + svg).
Khi sao chép XML sang slide khác, phải nhân bản luôn quan hệ (relationship)
ảnh, nếu không PowerPoint sẽ báo lỗi hỏng file. Hàm sao_hinh xử lý việc này.
"""
import copy
import os
import re

from pptx import Presentation
from pptx.util import Inches

GOC = os.path.dirname(os.path.abspath(__file__))
MAU = os.path.join(os.path.dirname(GOC), "Template.pptx")

# Khổ slide của mẫu
W_THAT, H_THAT = 20.0, 11.25

# Slide trong mẫu dùng làm nguồn khung
BIA = 0          # slide 1: trang bìa
MUC_LUC = 1      # slide 2: NỘI DUNG
CHUONG = 2       # slide 3: trang mở chương
NOI_DUNG = 5     # slide 6: trang nội dung (thanh tiêu đề + chân trang)
CAM_ON = 33      # slide 34: trang cảm ơn

_RID = re.compile(r'r:(?:embed|link)="([^"]+)"')


def _sao_quan_he(el, nguon_part, dich_part):
    """Nhân bản mọi ảnh nhúng trong el sang part đích, viết lại rId."""
    for rid in sorted(set(_RID.findall(el.xml))):
        try:
            anh = nguon_part.rels[rid].target_part
        except KeyError:
            continue
        rid_moi = dich_part.relate_to(anh, nguon_part.rels[rid].reltype)
        if rid_moi != rid:
            # đổi đúng thuộc tính r:embed / r:link, không đụng chuỗi khác
            for thuoc in ("embed", "link"):
                for n in el.iter():
                    for k in list(n.attrib):
                        if k.endswith("}" + thuoc) and n.attrib[k] == rid:
                            n.attrib[k] = rid_moi


def _co_chu(sh, mau_chu):
    """Shape (kể cả trong group) có chứa một trong các chuỗi cần bỏ?"""
    if not mau_chu:
        return False
    try:
        x = sh._element.xml
    except AttributeError:
        return False
    return any(m in x for m in mau_chu)


def sao_khung(nguon, dich, bo_anh=True, bo_chu=(), chi_lay=None):
    """Sao chép shape trang trí từ slide nguồn sang slide đích.

    bo_anh  : bỏ shape PICTURE (ảnh minh họa của mẫu, không phải trang trí)
    bo_chu  : bỏ shape có chứa các chuỗi này (chữ của bài mẫu)
    chi_lay : nếu đặt, chỉ lấy shape có tên trong danh sách này
    """
    for sh in nguon.shapes:
        if chi_lay is not None and sh.name not in chi_lay:
            continue
        if bo_anh and sh.shape_type == 13:          # PICTURE
            continue
        if sh.is_placeholder:                        # số trang: để layout lo
            continue
        if _co_chu(sh, bo_chu):
            continue
        el = copy.deepcopy(sh._element)
        _sao_quan_he(el, nguon.part, dich.part)
        dich.shapes._spTree.append(el)


def tim_theo_ten(slide, ten):
    """Tìm shape theo đúng tên (kể cả shape con trong group)."""
    def di(sh):
        if sh.name == ten:
            return sh
        if sh.shape_type == 6:
            for c in sh.shapes:
                r = di(c)
                if r is not None:
                    return r
        return None

    for sh in slide.shapes:
        r = di(sh)
        if r is not None:
            return r
    return None


def dat_chu_theo_ten(slide, ten_group, moi):
    """Đổi chữ của textbox nằm trong group có tên cho trước.

    Dùng cho thanh tiêu đề của mẫu: giữ nguyên mọi định dạng (font, cỡ,
    màu, vị trí) và chỉ thay nội dung chữ.
    """
    g = tim_theo_ten(slide, ten_group)
    if g is None:
        return False
    def di(sh):
        if sh.has_text_frame and sh.text_frame.paragraphs:
            p = sh.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = moi
                for r in p.runs[1:]:
                    r.text = ""
                for extra in sh.text_frame.paragraphs[1:]:
                    for r in extra.runs:
                        r.text = ""
                return True
        if sh.shape_type == 6:
            for c in sh.shapes:
                if di(c):
                    return True
        return False

    return di(g)


def tim_shape(slide, chuoi):
    """Tìm shape (hoặc shape con trong group) chứa đúng chuỗi chữ."""
    def di(sh):
        if sh.shape_type == 6:
            for c in sh.shapes:
                r = di(c)
                if r is not None:
                    return r
            return None
        if sh.has_text_frame and chuoi in sh.text_frame.text:
            return sh
        return None

    for sh in slide.shapes:
        r = di(sh)
        if r is not None:
            return r
    return None


def dat_chu(slide, cu, moi):
    """Đổi chữ trong shape của khung mẫu, giữ nguyên định dạng run đầu."""
    sh = tim_shape(slide, cu)
    if sh is None:
        return False
    tf = sh.text_frame
    p = tf.paragraphs[0]
    if not p.runs:
        return False
    p.runs[0].text = moi
    for r in p.runs[1:]:
        r.text = ""
    for extra in tf.paragraphs[1:]:
        for r in extra.runs:
            r.text = ""
    return True


def mo_mau():
    """Mở Template.pptx để đọc khung."""
    return Presentation(MAU)


_RID_ATTR = ("{http://schemas.openxmlformats.org/officeDocument/2006/"
             "relationships}id")


def xoa_slide_dau(pr, n):
    """Xóa n slide đầu (các slide gốc của mẫu), giữ master/layout/theme."""
    ids = pr.slides._sldIdLst
    for sid in list(ids)[:n]:
        pr.part.drop_rel(sid.get(_RID_ATTR))
        ids.remove(sid)


def xoa_het_slide(pr):
    """Xóa toàn bộ slide của mẫu, giữ lại master/layout/theme."""
    xoa_slide_dau(pr, len(pr.slides._sldIdLst))


def kiem_tra():
    pr = mo_mau()
    print("khổ: %.2f x %.2f in, %d slide, %d layout"
          % (pr.slide_width / Inches(1), pr.slide_height / Inches(1),
             len(pr.slides), len(pr.slide_masters[0].slide_layouts)))
    for i in (BIA, MUC_LUC, CHUONG, NOI_DUNG, CAM_ON):
        s = pr.slides[i]
        n = sum(1 for sh in s.shapes)
        print("  slide %2d: %d shape" % (i + 1, n))


if __name__ == "__main__":
    kiem_tra()
